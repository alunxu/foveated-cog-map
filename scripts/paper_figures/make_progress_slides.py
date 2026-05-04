"""Generate the polished CS503 Progress Report 2 slide deck.

9 slides, 3-min talk. Image-first layout, minimal but pointed text.
Findings slides: figure on top, 1-2 high-level takeaways below.
"""
from pathlib import Path
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ----- Constants ----- #
ROOT = Path(__file__).resolve().parent.parent.parent
FIG_DIR = ROOT / "docs/cs503_progress/fig"
MANU_FIG_DIR = ROOT / "docs/manuscript/fig"
OUT = ROOT / "docs/cs503_progress/CS503_Progress_Presentation.pptx"

# Palette (academic, restrained)
NAVY = RGBColor(0x1f, 0x3a, 0x68)
NAVY_LIGHT = RGBColor(0x4a, 0x5f, 0x82)
NAVY_DEEP = RGBColor(0x14, 0x29, 0x4a)
ACCENT = RGBColor(0xb0, 0x32, 0x29)        # restrained academic red
ACCENT_LIGHT = RGBColor(0xd5, 0x84, 0x7c)
TEXT = RGBColor(0x22, 0x22, 0x22)
LIGHTTEXT = RGBColor(0x6b, 0x6b, 0x6b)
BG_TINT = RGBColor(0xf8, 0xf9, 0xfb)
PANEL_BG = RGBColor(0xed, 0xf1, 0xf6)
PANEL_BG_WARM = RGBColor(0xf6, 0xee, 0xea)
DIV = RGBColor(0xc8, 0xd0, 0xdc)
HIGHLIGHT_BG = RGBColor(0xfe, 0xf6, 0xe0)
WHITE = RGBColor(0xff, 0xff, 0xff)


def make_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, left, top, w, h, color):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    fill(rect, color)
    return rect


def add_text(slide, text, left, top, w, h, *, size=18, bold=False,
             color=TEXT, italic=False, align=PP_ALIGN.LEFT,
             font="Times New Roman", anchor=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tx


def add_runs(slide, runs, left, top, w, h, *, align=PP_ALIGN.LEFT,
             font="Times New Roman", anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of (text, dict-of-formatting). Builds rich-text box."""
    tx = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, props in runs:
        r = p.add_run()
        r.text = text
        r.font.name = props.get("font", font)
        r.font.size = Pt(props.get("size", 16))
        r.font.bold = props.get("bold", False)
        r.font.italic = props.get("italic", False)
        r.font.color.rgb = props.get("color", TEXT)
    return tx


def add_top_bar(slide, title, *, kicker=None, accent=NAVY_DEEP):
    """Slide top bar — thin coloured strip + kicker + title + accent line."""
    add_rect(slide, 0, 0, 13.33, 0.18, accent)
    if kicker:
        add_text(slide, kicker, 0.6, 0.30, 12.1, 0.3,
                 size=11, color=ACCENT, bold=True, italic=False,
                 font="Helvetica")
        add_text(slide, title, 0.6, 0.55, 12.1, 0.7,
                 size=28, bold=True, color=NAVY)
    else:
        add_text(slide, title, 0.6, 0.36, 12.1, 0.8,
                 size=28, bold=True, color=NAVY)
    add_rect(slide, 0.6, 1.16, 1.4, 0.045, ACCENT)


def add_footer(slide, page, total=9, section=""):
    add_rect(slide, 0, 7.42, 13.33, 0.08, DIV)
    add_text(slide, f"CS-503  ·  Progress Report 2",
             0.6, 7.20, 6.0, 0.22,
             size=10, color=LIGHTTEXT, italic=True)
    add_text(slide, f"{section}    ·    {page} / {total}",
             7.0, 7.20, 5.7, 0.22,
             size=10, color=LIGHTTEXT, italic=True, align=PP_ALIGN.RIGHT)


def add_pdf_image(slide, pdf_path, left, top, *, width=None, height=None,
                  png_stem="_tmp_fig", dpi=200):
    out_stem = FIG_DIR / png_stem
    subprocess.run(["pdftoppm", "-png", "-r", str(dpi), str(pdf_path),
                    str(out_stem)], check=True)
    out_png = FIG_DIR / f"{png_stem}-1.png"
    if width:
        slide.shapes.add_picture(str(out_png), Inches(left), Inches(top),
                                  width=Inches(width))
    else:
        slide.shapes.add_picture(str(out_png), Inches(left), Inches(top),
                                  height=Inches(height))
    return out_png


def tint_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, BG_TINT)


def highlight_box(slide, text, left, top, w, h, *, color=NAVY,
                  bg=HIGHLIGHT_BG, size=15, bold=False):
    """Subtle highlight box for a key takeaway."""
    add_rect(slide, left, top, w, h, bg)
    add_rect(slide, left, top, 0.05, h, color)
    add_text(slide, text, left + 0.18, top + 0.05, w - 0.3, h - 0.1,
             size=size, color=TEXT, bold=bold, italic=False,
             anchor=MSO_ANCHOR.MIDDLE)


# ===============================================================
# Build slides
# ===============================================================
prs = make_prs()
blank = prs.slide_layouts[6]


# ===== Slide 1: Title ============================================ #
s = prs.slides.add_slide(blank)
tint_bg(s)
add_rect(s, 0, 0, 0.4, 7.5, NAVY_DEEP)

add_text(s, "CS-503 VISUAL INTELLIGENCE",
         0.9, 1.2, 11.4, 0.4,
         size=14, color=ACCENT, bold=True, italic=True, font="Helvetica")
add_text(s, "Progress Report 2",
         0.9, 1.55, 11.4, 0.5,
         size=18, color=LIGHTTEXT, italic=True)

# Stacked title with elegant kerning
add_text(s, "Sensor Structure Shapes",
         0.9, 2.5, 11.4, 0.85,
         size=44, bold=True, color=NAVY_DEEP)
add_text(s, "the Format of Cognitive Maps",
         0.9, 3.35, 11.4, 0.85,
         size=44, bold=True, color=NAVY_DEEP)
add_text(s, "in Navigation Agents",
         0.9, 4.20, 11.4, 0.85,
         size=44, bold=True, color=NAVY_DEEP)

# Author line
add_rect(s, 0.9, 5.45, 0.06, 0.25, ACCENT)
add_text(s, "Léo Bruneau    ·    Weilun Xu    ·    Rim Abkari    ·    Zyad Aoutir",
         1.05, 5.45, 11.4, 0.4, size=18, color=TEXT)

# Bottom strip
add_rect(s, 0.9, 6.45, 11.4, 0.04, ACCENT_LIGHT)
add_text(s, "Does visual-input structure shape the format of the cognitive map?",
         0.9, 6.55, 11.4, 0.5,
         size=16, italic=True, color=NAVY_LIGHT)


# ===== Slide 2: Motivation (cog-neuro precedent) ================== #
s = prs.slides.add_slide(blank)
tint_bg(s)
add_top_bar(s, "Different sensor ecologies, different cognitive maps",
            kicker="MOTIVATION")

photo_row_y = 1.55
photo_row_h = 2.7
photo_w = 3.7
gap = 0.4
photo_xs = [0.6, 0.6 + photo_w + gap, 0.6 + 2 * (photo_w + gap)]

panels = [
    ("photo_bat.jpg", "Echolocating bat", "no eyes  ·  sonar map"),
    ("photo_mole.jpg", "Blind rodent", "tactile  ·  path integration"),
    ("photo_eye.jpg", "Foveated primate", "fovea + blurred periphery"),
]
for (fname, label, tagline), x in zip(panels, photo_xs):
    s.shapes.add_picture(str(FIG_DIR / fname),
                         Inches(x), Inches(photo_row_y),
                         width=Inches(photo_w), height=Inches(photo_row_h))
    add_rect(s, x, photo_row_y + photo_row_h, photo_w, 0.92, PANEL_BG)
    add_rect(s, x, photo_row_y + photo_row_h, photo_w, 0.04, NAVY)
    add_text(s, label, x + 0.1, photo_row_y + photo_row_h + 0.10,
             photo_w - 0.2, 0.4, size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, tagline, x + 0.1, photo_row_y + photo_row_h + 0.50,
             photo_w - 0.2, 0.35, size=12, italic=True, color=LIGHTTEXT,
             align=PP_ALIGN.CENTER)

synth_y = photo_row_y + photo_row_h + 1.10
add_text(s, "→ qualitatively different hippocampal coding strategies, "
         "not just different accuracy",
         0.6, synth_y, 12.1, 0.4,
         size=15, italic=True, color=NAVY_LIGHT, align=PP_ALIGN.CENTER)

add_rect(s, 0.6, synth_y + 0.6, 12.1, 0.65, NAVY_DEEP)
add_text(s, "Does varying only the visual sensor change the cognitive map's format in deep-RL agents?",
         0.6, synth_y + 0.68, 12.1, 0.5,
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_footer(s, page=2, section="Motivation")


# ===== Slide 3: Setup ============================================= #
s = prs.slides.add_slide(blank)
tint_bg(s)
add_top_bar(s, "5 conditions on one bandwidth axis",
            kicker="SETUP")

add_pdf_image(s, MANU_FIG_DIR / "fig_pipeline_overview_tikz.pdf",
              0.6, 1.45, width=10.0, png_stem="_pipeline", dpi=180)

panel_left = 11.0
panel_w = 1.95
add_rect(s, panel_left, 1.45, panel_w, 5.5, PANEL_BG)
add_rect(s, panel_left, 1.45, panel_w, 0.05, NAVY)

add_text(s, "Three lenses", panel_left + 0.15, 1.58, panel_w - 0.3, 0.35,
         size=14, bold=True, color=NAVY)
add_text(s, "on the frozen ℎ₂",
         panel_left + 0.15, 1.92, panel_w - 0.3, 0.3,
         size=11, italic=True, color=LIGHTTEXT)

# Three small bullets
for i, (title, sub) in enumerate([
    ("Decodability", "what does ℎ₂ encode?"),
    ("Intrinsic structure", "do conditions share a code?"),
    ("Causal intervention", "does the policy rely on ℎ₂?"),
]):
    y = 2.4 + i * 1.45
    add_rect(s, panel_left + 0.15, y, 0.05, 0.55, ACCENT)
    add_text(s, title, panel_left + 0.28, y, panel_w - 0.4, 0.3,
             size=11, bold=True, color=NAVY)
    add_text(s, sub, panel_left + 0.28, y + 0.30, panel_w - 0.4, 0.6,
             size=10, italic=True, color=LIGHTTEXT)

add_text(s, "Same architecture, training, sensors. "
         "All conditions reach 96–99 % success — no skill gap.",
         0.6, 6.6, 12.1, 0.3,
         size=13, italic=True, color=NAVY_LIGHT, align=PP_ALIGN.CENTER)

add_footer(s, page=3, section="Setup")


# ===== Slide 4: Project rescope =================================== #
s = prs.slides.add_slide(blank)
tint_bg(s)
add_top_bar(s, "Project rescope: from foveation to encoder bandwidth",
            kicker="WHAT CHANGED SINCE THE PROPOSAL")

# Left: H1
col_w = 6.05
add_rect(s, 0.6, 1.55, col_w, 5.4, PANEL_BG)
add_rect(s, 0.6, 1.55, 0.08, 5.4, ACCENT)
add_text(s, "H1 reframed", 0.85, 1.75, col_w - 0.3, 0.5,
         size=22, bold=True, color=ACCENT)
add_text(s, "Encoder bandwidth, not foveation.",
         0.85, 2.25, col_w - 0.3, 0.4,
         size=14, italic=True, color=NAVY_LIGHT)
add_runs(s, [
    ("Unified-hyperparameter training removes apparent ",
     {"size": 15}),
    ("foveated specialness",
     {"size": 15, "bold": True, "color": NAVY}),
    (".  Foveated-logpolar (2×2) control lands where the ",
     {"size": 15}),
    ("encoder-bandwidth account predicts",
     {"size": 15, "bold": True, "color": NAVY}),
    (".",
     {"size": 15}),
], 0.85, 3.0, col_w - 0.3, 4.0, line_spacing=1.4)

# Right: H3
add_rect(s, 7.0, 1.55, col_w, 5.4, PANEL_BG_WARM)
add_rect(s, 7.0, 1.55, 0.08, 5.4, ACCENT)
add_text(s, "H3 de-scoped",
         7.25, 1.75, col_w - 0.3, 0.5, size=22, bold=True, color=ACCENT)
add_text(s, "Cog-neuro toolkit imported instead.",
         7.25, 2.25, col_w - 0.3, 0.4,
         size=14, italic=True, color=NAVY_LIGHT)
add_runs(s, [
    ("Learned-gaze MLP collapsed under navigation reward.  Redirected to ",
     {"size": 15}),
    ("Skaggs · SR predictive horizon · Procrustes · subspace divergence · excursion-forgetting",
     {"size": 15, "bold": True, "color": NAVY}),
    (".",
     {"size": 15}),
], 7.25, 3.0, col_w - 0.3, 4.0, line_spacing=1.4)

add_footer(s, page=4, section="Rescope")


# ===== Helper for finding slides (figure on top, takeaways below) = #
def finding_slide(prs, title, kicker, fig_path, png_stem,
                   takeaway_main, takeaway_sub, page, section):
    s = prs.slides.add_slide(blank)
    tint_bg(s)
    add_top_bar(s, title, kicker=kicker)

    # Figure: full-width, top half
    add_pdf_image(s, fig_path, 1.4, 1.4,
                  width=10.5, png_stem=png_stem, dpi=200)

    # Takeaway band: highlight-style
    main_y = 5.55
    add_rect(s, 0.6, main_y, 12.1, 0.95, HIGHLIGHT_BG)
    add_rect(s, 0.6, main_y, 0.06, 0.95, NAVY)
    add_text(s, takeaway_main, 0.85, main_y + 0.13, 11.6, 0.7,
             size=18, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    # sub-takeaway in italic, small
    add_text(s, takeaway_sub, 0.6, main_y + 1.05, 12.1, 0.4,
             size=13, italic=True, color=NAVY_LIGHT,
             align=PP_ALIGN.CENTER)
    add_footer(s, page=page, section=section)
    return s


# ===== Slide 5: Finding 1 ========================================= #
finding_slide(prs,
    title="Encoder bandwidth shapes the linear axis",
    kicker="FINDING 1 / 3   ·   MAGNITUDE AXIS",
    fig_path=FIG_DIR / "fig1_capacity_allocation.pdf",
    png_stem="_fig1",
    takeaway_main="Linear GPS R² falls monotonically with encoder bandwidth — yet MLP recovers GPS in every condition.",
    takeaway_sub="Position is retained, just not as the linearly-readable axis. Format differs.",
    page=5, section="Magnitude")


# ===== Slide 6: Finding 2 ========================================= #
finding_slide(prs,
    title="Allocation emerges across training",
    kicker="FINDING 2 / 3   ·   TRAINING-TIME PROCESS",
    fig_path=FIG_DIR / "fig3_substitution_dynamics.pdf",
    png_stem="_fig2",
    takeaway_main="All conditions start with a high linear GPS code; rich-encoder ones decay as the visual route consolidates.",
    takeaway_sub="Decay rate ordered by encoder informativeness — capacity allocation is a training-time process.",
    page=6, section="Training-time")


# ===== Slide 7: Finding 3 ========================================= #
finding_slide(prs,
    title="Bottleneck integrates; rich encoder reacts",
    kicker="FINDING 3 / 3   ·   FORMAT GENERALISES",
    fig_path=FIG_DIR / "fig2_lagk_stability.pdf",
    png_stem="_fig3",
    takeaway_main="Bottleneck stable past-position decoding across lags; rich-encoder degrades — generalises GPS → compass.",
    takeaway_sub="It is about format, not GPS specifically.",
    page=7, section="Format generalises")


# ===== Slide 8: Synthesis (3-axis) =============================== #
s = prs.slides.add_slide(blank)
tint_bg(s)
add_top_bar(s, "Three regimes of cognitive-map format",
            kicker="SYNTHESIS")

# Synthesis figure (left, big)
add_pdf_image(s, MANU_FIG_DIR / "fig6_synthesis_2axes.pdf",
              0.6, 1.5, width=8.4, png_stem="_syn", dpi=180)

# Right-side narrative
panel_left = 9.3
panel_w = 3.5
add_rect(s, panel_left, 1.5, panel_w, 5.5, PANEL_BG)
add_rect(s, panel_left, 1.5, panel_w, 0.05, NAVY)

add_text(s, "Three regimes",
         panel_left + 0.2, 1.65, panel_w - 0.4, 0.4,
         size=16, bold=True, color=NAVY)

# Regime descriptions
for i, (title, desc, color) in enumerate([
    ("Bottleneck", "Linear GPS code · format-shared",
     NAVY),
    ("Rich-encoder", "No linear code · format-isolated",
     ACCENT),
    ("No-vision", "Linear code · format-isolated",
     NAVY_LIGHT),
]):
    y = 2.2 + i * 1.0
    add_rect(s, panel_left + 0.2, y, 0.05, 0.7, color)
    add_text(s, title,
             panel_left + 0.32, y, panel_w - 0.5, 0.3,
             size=13, bold=True, color=color)
    add_text(s, desc,
             panel_left + 0.32, y + 0.30, panel_w - 0.5, 0.4,
             size=11, italic=True, color=LIGHTTEXT)

# Bottom: scientific bet
add_text(s, "Long-term scientific bet",
         panel_left + 0.2, 5.4, panel_w - 0.4, 0.3,
         size=12, bold=True, color=ACCENT)
add_text(s, "Does the bandwidth principle predict how real "
         "hippocampal circuits differ across species' sensor ecologies?",
         panel_left + 0.2, 5.7, panel_w - 0.4, 1.2,
         size=11, italic=True, color=NAVY_LIGHT)

add_footer(s, page=8, section="Synthesis")


# ===== Slide 9: Open issues + Implications ====================== #
s = prs.slides.add_slide(blank)
tint_bg(s)
add_top_bar(s, "Open issues  ·  Implications",
            kicker="WHAT'S NEXT  ·  WHAT IT MEANS")

# Left: action items — icon-card layout
col_w = 6.05
add_rect(s, 0.6, 1.55, col_w, 5.4, PANEL_BG)
add_rect(s, 0.6, 1.55, 0.08, 5.4, NAVY)
add_text(s, "Toward the final report",
         0.85, 1.75, col_w - 0.3, 0.5,
         size=22, bold=True, color=NAVY)
add_text(s, "Three batches of pending work.",
         0.85, 2.25, col_w - 0.3, 0.4,
         size=13, italic=True, color=LIGHTTEXT)

action_items = [
    ("icon_action_1_bandwidth.png",
     "Bandwidth refinements",
     "blind cross-training",
     "coarse transplant",
     "two-seed validation"),
    ("icon_action_2_cogneuro.png",
     "Cog-neuro on ℎ₂",
     "occupancy decoder",
     "MINE-PH · SR eigenbasis",
     "Sussillo–Barak fixed-points"),
    ("icon_action_3_scope.png",
     "Scope tests",
     "foveation strength sweep",
     "transformer falsification",
     ""),
]

card_top = 2.85
card_h = 1.36
icon_size = 1.05
icon_left = 0.95
text_left = icon_left + icon_size + 0.18

for i, (icon, title, *items) in enumerate(action_items):
    y = card_top + i * card_h
    # Icon
    s.shapes.add_picture(str(FIG_DIR / icon),
                         Inches(icon_left), Inches(y + 0.08),
                         width=Inches(icon_size), height=Inches(icon_size))
    # Vertical accent
    add_rect(s, text_left - 0.04, y + 0.08, 0.04, icon_size, ACCENT)
    # Title
    add_text(s, title, text_left + 0.08, y + 0.08,
             col_w - text_left - 0.2, 0.36,
             size=14, bold=True, color=NAVY)
    # Items as comma-separated mini-list
    item_text = "  ·  ".join([it for it in items if it])
    add_text(s, item_text, text_left + 0.08, y + 0.45,
             col_w - text_left - 0.2, 0.7,
             size=10, italic=True, color=LIGHTTEXT, anchor=MSO_ANCHOR.TOP)

# Right: implications
add_rect(s, 7.0, 1.55, col_w, 5.4, PANEL_BG_WARM)
add_rect(s, 7.0, 1.55, 0.08, 5.4, ACCENT)
add_text(s, "Implications", 7.25, 1.75, col_w - 0.3, 0.5,
         size=22, bold=True, color=ACCENT)
add_text(s, "If the principle survives the remaining tests …",
         7.25, 2.25, col_w - 0.3, 0.4,
         size=13, italic=True, color=LIGHTTEXT)

for i, (k, v) in enumerate([
    ("Encoder bandwidth as design lever",
     "restricting the encoder may recruit integration-style memory"),
    ("Architecture-agnostic in principle",
     "VLMs on rich foundation encoders may show the same crowd-out"),
    ("Methodology porting",
     "linear / lag-k / place-cell lens → cog-neuro tools for deep RL"),
]):
    y = 2.95 + i * 1.25
    add_rect(s, 7.25, y, 0.06, 0.85, NAVY)
    add_text(s, k, 7.4, y, col_w - 0.5, 0.35,
             size=14, bold=True, color=ACCENT)
    add_text(s, v, 7.4, y + 0.35, col_w - 0.5, 0.55,
             size=11, italic=True, color=LIGHTTEXT)

add_footer(s, page=9, section="Open issues  ·  Implications")


# ===== Save ====================================================== #
prs.save(str(OUT))
print(f"wrote {OUT}")

# Cleanup intermediate PNGs
for stem in ["_pipeline-1.png", "_fig1-1.png", "_fig2-1.png", "_fig3-1.png",
             "_syn-1.png"]:
    p = FIG_DIR / stem
    if p.exists():
        p.unlink()
